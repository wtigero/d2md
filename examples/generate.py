#!/usr/bin/env python3
"""Generate small, reproducible examples for every format d2md accepts.

The output is intentionally generated rather than committed as binary blobs:
the recipe is reviewable, free of client data, and useful from
both macOS/Linux shells and Windows PowerShell.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
import struct
import zipfile


DOCUMENTS: list[dict[str, str]] = []
FIXED_ZIP_TIME = (2026, 1, 1, 0, 0, 0)


def record(
    path: Path,
    expected_text: str,
    description: str,
    capability: str = "base",
) -> None:
    DOCUMENTS.append(
        {
            "file": path.name,
            "extension": path.suffix.lower(),
            "expected_text": expected_text,
            "description": description,
            "capability": capability,
        }
    )


def draw_ocr_marker(image, text: str) -> None:
    """Draw a large anti-aliased marker on each supported smoke-test OS."""
    import os

    from PIL import ImageDraw, ImageFont

    windows = Path(os.environ.get("WINDIR", "C:/Windows"))
    candidates = (
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        windows / "Fonts" / "arial.ttf",
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/dejavu/DejaVuSans.ttf"),
    )
    font = next(
        (
            ImageFont.truetype(str(candidate), 64)
            for candidate in candidates
            if candidate.is_file()
        ),
        None,
    )
    if font is None:
        try:
            font = ImageFont.load_default(size=64)
        except TypeError:
            font = ImageFont.load_default()
    ImageDraw.Draw(image).text((80, 180), text, fill="black", font=font)


def write_text(
    output: Path,
    filename: str,
    text: str,
    description: str,
    encoding: str = "utf-8",
    expected_text: str | None = None,
) -> None:
    path = output / filename
    path.write_text(text, encoding=encoding)
    record(path, expected_text or text.splitlines()[0], description)


def write_plain_examples(output: Path) -> None:
    write_text(
        output,
        "plain-utf8.txt",
        "D2MD EXAMPLE TXT\nทดสอบภาษาไทย น้ำจำกัด ระบบอ่านไฟล์ทดสอบ\n",
        "UTF-8 plain text with Thai",
    )
    write_text(
        output,
        "plain-cp874.txt",
        "D2MD EXAMPLE CP874\nทดสอบภาษาไทย น้ำจำกัด ระบบอ่านไฟล์ทดสอบ\n",
        "Legacy Thai CP874 text",
        encoding="cp874",
    )
    write_text(
        output,
        "markdown.md",
        "# D2MD EXAMPLE MARKDOWN\n\nA small Markdown document.\n",
        "Markdown",
    )
    write_text(
        output,
        "table.csv",
        "D2MD EXAMPLE CSV,Value\nEnglish,42\nThai,ทดสอบ\n",
        "Comma-separated values",
    )
    write_text(
        output,
        "data-json.json",
        json.dumps(
            {"title": "D2MD EXAMPLE JSON", "thai": "ทดสอบ", "value": 42},
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        "JSON",
        expected_text="D2MD EXAMPLE JSON",
    )
    write_text(
        output,
        "data-xml.xml",
        "<?xml version=\"1.0\" encoding=\"UTF-8\"?>\n"
        "<document><title>D2MD EXAMPLE XML</title><thai>ทดสอบ</thai></document>\n",
        "XML",
        expected_text="D2MD EXAMPLE XML",
    )
    write_text(
        output,
        "data-yaml.yaml",
        "title: D2MD EXAMPLE YAML\nthai: ทดสอบ\nvalue: 42\n",
        "YAML",
    )
    write_text(
        output,
        "data-yml.yml",
        "title: D2MD EXAMPLE YML\nthai: ทดสอบ\nvalue: 42\n",
        "YML alias",
    )


def write_html_examples(output: Path) -> None:
    body = (
        "<!doctype html><html><head><meta charset=\"utf-8\">"
        "<title>D2MD EXAMPLE HTML</title></head><body>"
        "<h1>D2MD EXAMPLE HTML</h1><p>Thai: ทดสอบภาษาไทย</p>"
        "<table><tr><th>Item</th><th>Value</th></tr>"
        "<tr><td>Example</td><td>42</td></tr></table></body></html>"
    )
    write_text(
        output,
        "web-page.html",
        body,
        "HTML page",
        expected_text="D2MD EXAMPLE HTML",
    )
    write_text(
        output,
        "web-page-legacy.htm",
        body.replace("EXAMPLE HTML", "EXAMPLE HTM"),
        "HTM alias",
        expected_text="D2MD EXAMPLE HTM",
    )


def write_openxml_examples(output: Path) -> None:
    try:
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise SystemExit(
            "OpenXML example dependencies are missing; install with: "
            "python -m pip install -e '.[dev]'"
        ) from exc

    docx_path = output / "word-document.docx"
    document = Document()
    document.add_heading("D2MD EXAMPLE DOCX", level=1)
    document.add_paragraph("Thai: ทดสอบภาษาไทย น้ำจำกัด")
    document.add_table(rows=2, cols=2).cell(0, 0).text = "Example table"
    document.save(docx_path)
    record(docx_path, "D2MD EXAMPLE DOCX", "Word document")

    xlsx_path = output / "spreadsheet.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Example"
    sheet.append(["D2MD EXAMPLE XLSX", "Value"])
    sheet.append(["Thai: ทดสอบ", 42])
    workbook.save(xlsx_path)
    record(xlsx_path, "D2MD EXAMPLE XLSX", "Excel workbook")

    pptx_path = output / "presentation.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text_frame.text = "D2MD EXAMPLE PPTX\nThai: ทดสอบภาษาไทย"
    presentation.save(pptx_path)
    record(pptx_path, "D2MD EXAMPLE PPTX", "PowerPoint presentation")


def write_xls_example(output: Path) -> None:
    try:
        import xlwt
    except ImportError as exc:
        raise SystemExit(
            "The .xls generator needs xlwt; install with: "
            "python -m pip install -e '.[dev]'"
        ) from exc

    path = output / "spreadsheet-legacy.xls"
    workbook = xlwt.Workbook()
    sheet = workbook.add_sheet("Example")
    sheet.write(0, 0, "D2MD EXAMPLE XLS")
    sheet.write(0, 1, "Value")
    sheet.write(1, 0, "Thai: ทดสอบ")
    sheet.write(1, 1, 42)
    workbook.save(str(path))
    record(path, "D2MD EXAMPLE XLS", "Legacy BIFF Excel workbook")


def zip_write(archive: zipfile.ZipFile, name: str, content: str) -> None:
    info = zipfile.ZipInfo(name, FIXED_ZIP_TIME)
    info.compress_type = zipfile.ZIP_DEFLATED
    archive.writestr(info, content.encode("utf-8"))


def write_epub_example(output: Path) -> None:
    path = output / "ebook.epub"
    with zipfile.ZipFile(path, "w") as archive:
        mimetype = zipfile.ZipInfo("mimetype", FIXED_ZIP_TIME)
        mimetype.compress_type = zipfile.ZIP_STORED
        archive.writestr(mimetype, b"application/epub+zip")
        zip_write(
            archive,
            "META-INF/container.xml",
            "<?xml version=\"1.0\"?>"
            "<container version=\"1.0\" "
            "xmlns=\"urn:oasis:names:tc:opendocument:xmlns:container\">"
            "<rootfiles><rootfile full-path=\"OEBPS/content.opf\" "
            "media-type=\"application/oebps-package+xml\"/></rootfiles></container>",
        )
        zip_write(
            archive,
            "OEBPS/content.opf",
            "<?xml version=\"1.0\" encoding=\"UTF-8\"?>"
            "<package xmlns=\"http://www.idpf.org/2007/opf\" version=\"3.0\" "
            "unique-identifier=\"book-id\"><metadata "
            "xmlns:dc=\"http://purl.org/dc/elements/1.1/\">"
            "<dc:identifier id=\"book-id\">d2md-example</dc:identifier>"
            "<dc:title>D2MD EXAMPLE EPUB</dc:title>"
            "<dc:language>en</dc:language></metadata>"
            "<manifest><item id=\"chapter\" href=\"chapter.xhtml\" "
            "media-type=\"application/xhtml+xml\"/></manifest>"
            "<spine><itemref idref=\"chapter\"/></spine></package>",
        )
        zip_write(
            archive,
            "OEBPS/chapter.xhtml",
            "<?xml version=\"1.0\" encoding=\"UTF-8\"?>"
            "<html xmlns=\"http://www.w3.org/1999/xhtml\"><head>"
            "<title>D2MD EXAMPLE EPUB</title></head><body>"
            "<h1>D2MD EXAMPLE EPUB</h1><p>Thai: ทดสอบภาษาไทย</p>"
            "</body></html>",
        )
    record(path, "D2MD EXAMPLE EPUB", "EPUB e-book")


def directory_entry(
    name: str,
    object_type: int,
    start_sector: int,
    size: int,
    left: int = 0xFFFFFFFF,
    right: int = 0xFFFFFFFF,
    child: int = 0xFFFFFFFF,
) -> bytes:
    encoded_name = (name + "\0").encode("utf-16-le")
    if len(encoded_name) > 64:
        raise ValueError(f"OLE directory name is too long: {name}")
    entry = bytearray(128)
    entry[: len(encoded_name)] = encoded_name
    struct.pack_into("<HBBIII", entry, 64, len(encoded_name), object_type, 1, left, right, child)
    struct.pack_into("<I", entry, 116, start_sector)
    struct.pack_into("<Q", entry, 120, size)
    return bytes(entry)


def write_msg_example(output: Path) -> None:
    """Write a minimal valid Outlook MSG compound file without private data."""
    free_sector = 0xFFFFFFFF
    end_of_chain = 0xFFFFFFFE
    fat_sector = 0xFFFFFFFD
    fat_index = 34

    header = bytearray(512)
    header[:8] = bytes.fromhex("D0CF11E0A1B11AE1")
    struct.pack_into("<HHHH", header, 24, 0x003E, 0x0003, 0xFFFE, 9)
    struct.pack_into("<H", header, 32, 6)
    struct.pack_into("<IIIIIIIII", header, 40, 0, 1, 0, 0, 4096, end_of_chain, 0, end_of_chain, 0)
    struct.pack_into("<I", header, 76, fat_index)
    for offset in range(80, 512, 4):
        struct.pack_into("<I", header, offset, free_sector)

    streams = [
        ("__substg1.0_0037001F", "D2MD EXAMPLE MSG"),
        ("__substg1.0_0C1F001F", "sender@example.invalid"),
        ("__substg1.0_0E04001F", "reader@example.invalid"),
        (
            "__substg1.0_1000001F",
            "D2MD EXAMPLE MSG body. Thai: ทดสอบภาษาไทย น้ำจำกัด",
        ),
    ]
    starts = [2, 10, 18, 26]
    entries = [
        directory_entry("Root Entry", 5, end_of_chain, 0, child=2),
        directory_entry(streams[0][0], 2, starts[0], 4096),
        directory_entry(streams[1][0], 2, starts[1], 4096, left=1, right=3),
        directory_entry(streams[2][0], 2, starts[2], 4096, right=4),
        directory_entry(streams[3][0], 2, starts[3], 4096),
    ]
    directory = b"".join(entries).ljust(1024, b"\0")

    data = bytearray()
    for _name, value in streams:
        encoded = value.encode("utf-16-le")
        padding = " ".encode("utf-16-le") * ((4096 - len(encoded)) // 2)
        data.extend((encoded + padding).ljust(4096, b" "))

    fat = [free_sector] * 128
    fat[0], fat[1] = 1, end_of_chain
    for start in starts:
        for sector in range(start, start + 7):
            fat[sector] = sector + 1
        fat[start + 7] = end_of_chain
    fat[fat_index] = fat_sector
    fat_bytes = struct.pack("<128I", *fat)

    path = output / "email.msg"
    path.write_bytes(bytes(header) + directory + bytes(data) + fat_bytes)
    record(path, "D2MD EXAMPLE MSG", "Outlook MSG email")


def write_pdf_examples(output: Path) -> None:
    try:
        from PIL import Image
        from reportlab.pdfgen import canvas
    except ImportError as exc:
        raise SystemExit(
            "PDF example dependencies are missing; install with: "
            "python -m pip install -e '.[dev]'"
        ) from exc

    digital = output / "pdf-born-digital.pdf"
    pdf = canvas.Canvas(str(digital))
    pdf.drawString(72, 720, "D2MD EXAMPLE PDF BORN DIGITAL")
    pdf.drawString(72, 700, "This page has a searchable text layer and the value 42.")
    pdf.save()
    record(digital, "D2MD EXAMPLE PDF BORN DIGITAL", "Born-digital PDF")

    scan_image = Image.new("RGB", (1800, 500), "white")
    draw_ocr_marker(scan_image, "D2MD EXAMPLE PDF SCAN 42")
    scanned = output / "pdf-scanned.pdf"
    scan_image.save(scanned, "PDF", resolution=150.0)
    record(
        scanned,
        "D2MD EXAMPLE PDF SCAN",
        "Image-only scanned PDF",
        capability="ocr",
    )


def write_image_examples(output: Path) -> None:
    try:
        from PIL import Image
    except ImportError as exc:
        raise SystemExit(
            "Image example dependencies are missing; install with: "
            "python -m pip install -e '.[dev]'"
        ) from exc

    formats = [
        ("image-png.png", "PNG", "D2MD EXAMPLE IMAGE PNG"),
        ("image-jpg.jpg", "JPEG", "D2MD EXAMPLE IMAGE JPG"),
        ("image-jpeg.jpeg", "JPEG", "D2MD EXAMPLE IMAGE JPEG"),
        ("image-tiff.tiff", "TIFF", "D2MD EXAMPLE IMAGE TIFF"),
        ("image-tif.tif", "TIFF", "D2MD EXAMPLE IMAGE TIF"),
        ("image-bmp.bmp", "BMP", "D2MD EXAMPLE IMAGE BMP"),
        ("image-webp.webp", "WEBP", "D2MD EXAMPLE IMAGE WEBP"),
    ]
    for filename, image_format, text in formats:
        image = Image.new("RGB", (1800, 500), "white")
        draw_ocr_marker(image, text + " 42")
        path = output / filename
        image.save(path, image_format)
        record(path, text, f"{image_format} image", capability="ocr")


def write_manifest(output: Path) -> None:
    manifest = {
        "generated": True,
        "contains_private_data": False,
        "documents": sorted(DOCUMENTS, key=lambda entry: entry["file"]),
    }
    (output / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def generate(output: Path) -> None:
    output.mkdir(parents=True, exist_ok=True)
    DOCUMENTS.clear()
    write_plain_examples(output)
    write_html_examples(output)
    write_openxml_examples(output)
    write_xls_example(output)
    write_epub_example(output)
    write_msg_example(output)
    write_pdf_examples(output)
    write_image_examples(output)
    write_manifest(output)


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Generate synthetic sample documents for every d2md format."
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path(__file__).parent / "generated",
        help="destination directory (default: examples/generated)",
    )
    args = parser.parse_args()
    generate(args.output.expanduser())
    print(f"generated {len(DOCUMENTS)} documents in {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
